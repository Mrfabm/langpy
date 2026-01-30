import os
from pathlib import Path

# Text-based
TEXT = "This is a test text file.\nIt has multiple lines."
MARKDOWN = "# Test Document\n\nThis is **markdown** content."
CSV = "Name,Age,City\nJohn,25,NY\nJane,30,LA"
TSV = "Name\tAge\tCity\nJohn\t25\tNY\nJane\t30\tLA"
HTML = "<html><body><h1>Test</h1><p>Content</p></body></html>"

# Output directory
outdir = Path(__file__).parent / "test_files"
outdir.mkdir(exist_ok=True)

# Write text files
def write_text():
    (outdir / "test.txt").write_text(TEXT)
    (outdir / "test.md").write_text(MARKDOWN)
    (outdir / "test.csv").write_text(CSV)
    (outdir / "test.tsv").write_text(TSV)
    (outdir / "test.html").write_text(HTML)

# Write docx
def write_docx():
    from docx import Document
    doc = Document()
    doc.add_heading("Test Document", 0)
    doc.add_paragraph("This is a test Word document.")
    doc.save(outdir / "test.docx")

# Write pptx
def write_pptx():
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Test Presentation"
    prs.save(outdir / "test.pptx")

# Write xlsx
def write_xlsx():
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.append(["Name", "Age", "City"])
    ws.append(["John", 25, "NY"])
    ws.append(["Jane", 30, "LA"])
    wb.save(outdir / "test.xlsx")

# Write pdf
def write_pdf():
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    c = canvas.Canvas(str(outdir / "test.pdf"), pagesize=letter)
    c.drawString(100, 750, "This is a test PDF document.")
    c.save()

# Write images
def write_images():
    from PIL import Image
    img = Image.new("RGB", (100, 100), color=(73, 109, 137))
    img.save(outdir / "test.jpg", "JPEG")
    img.save(outdir / "test.jpeg", "JPEG")
    img.save(outdir / "test.png", "PNG")
    img = Image.new("RGB", (100, 100), color=(255, 255, 0))
    img.save(outdir / "test.tiff", "TIFF")

if __name__ == "__main__":
    write_text()
    write_docx()
    write_pptx()
    write_xlsx()
    write_pdf()
    write_images()
    print(f"Test files written to {outdir}") 