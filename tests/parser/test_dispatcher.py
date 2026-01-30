"""
Test Dispatcher - Comprehensive testing for the new parser dispatcher.

Tests all supported formats, error handling, and dependency checks.
"""

import sys
import os
from pathlib import Path
import time
import io

# Add project root to path
project_root = Path(__file__).parent.parent.parent
sys.path.insert(0, str(project_root))

from parser.models import ParseRequest, ParserOptions
from parser.docling_parser_full_copy import parse_sync_dispatcher, SUPPORTED_FORMATS


def create_test_files():
    """Create test files for each supported format."""
    test_files = {}
    
    # Text files
    test_files['test.txt'] = b"This is a test text file.\nIt has multiple lines.\nAnd some content."
    test_files['test.md'] = b"# Test Markdown\n\nThis is a **test** markdown file.\n\n- Item 1\n- Item 2"
    test_files['test.csv'] = b"Name,Age,City\nJohn,30,New York\nJane,25,Los Angeles"
    
    # HTML file
    test_files['test.html'] = b"""
    <html>
    <head><title>Test HTML</title></head>
    <body>
        <h1>Test HTML Document</h1>
        <p>This is a test HTML file with some content.</p>
        <ul>
            <li>Item 1</li>
            <li>Item 2</li>
        </ul>
    </body>
    </html>
    """
    
    # Create XLSX file using openpyxl
    try:
        import pandas as pd
        df = pd.DataFrame({
            'Name': ['John', 'Jane', 'Bob'],
            'Age': [30, 25, 35],
            'City': ['New York', 'Los Angeles', 'Chicago']
        })
        buffer = io.BytesIO()
        df.to_excel(buffer, index=False)
        test_files['test.xlsx'] = buffer.getvalue()
        print("✅ Created test.xlsx")
    except ImportError:
        print("⚠️  pandas/openpyxl not available - skipping XLSX test")
    
    # Create DOCX file using python-docx
    try:
        from docx import Document
        doc = Document()
        doc.add_heading('Test Document', 0)
        doc.add_paragraph('This is a test Word document.')
        doc.add_paragraph('It contains multiple paragraphs.')
        buffer = io.BytesIO()
        doc.save(buffer)
        test_files['test.docx'] = buffer.getvalue()
        print("✅ Created test.docx")
    except ImportError:
        print("⚠️  python-docx not available - skipping DOCX test")
    
    # Create PPTX file using python-pptx
    try:
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = "Test Presentation"
        content = slide.placeholders[1]
        content.text = "This is a test PowerPoint presentation."
        buffer = io.BytesIO()
        prs.save(buffer)
        test_files['test.pptx'] = buffer.getvalue()
        print("✅ Created test.pptx")
    except ImportError:
        print("⚠️  python-pptx not available - skipping PPTX test")
    
    # Create image files using Pillow
    try:
        from PIL import Image, ImageDraw, ImageFont
        # Create a simple image with text
        img = Image.new('RGB', (400, 200), color='white')
        draw = ImageDraw.Draw(img)
        draw.text((50, 50), "Test Image", fill='black')
        draw.text((50, 100), "This is a test image file", fill='black')
        
        # Save as JPEG
        buffer = io.BytesIO()
        img.save(buffer, format='JPEG')
        test_files['test.jpg'] = buffer.getvalue()
        print("✅ Created test.jpg")
        
        # Save as PNG
        buffer = io.BytesIO()
        img.save(buffer, format='PNG')
        test_files['test.png'] = buffer.getvalue()
        print("✅ Created test.png")
    except ImportError:
        print("⚠️  Pillow not available - skipping image tests")
    
    # Create PDF file using PyPDF2
    try:
        import PyPDF2
        # Create a simple PDF
        writer = PyPDF2.PdfWriter()
        page = PyPDF2.PageObject.create_blank_page(width=612, height=792)
        writer.add_page(page)
        buffer = io.BytesIO()
        writer.write(buffer)
        test_files['test.pdf'] = buffer.getvalue()
        print("✅ Created test.pdf")
    except ImportError:
        print("⚠️  PyPDF2 not available - skipping PDF test")
    
    return test_files


def test_dispatcher():
    """Test the dispatcher with all supported formats."""
    print("🚀 Testing Parser Dispatcher")
    print("=" * 50)
    
    # Create test files
    print("\n📁 Creating test files...")
    test_files = create_test_files()
    
    if not test_files:
        print("❌ No test files could be created. Please check dependencies.")
        return
    
    print(f"\n✅ Created {len(test_files)} test files")
    
    # Test each file
    results = {}
    errors = {}
    
    for filename, content in test_files.items():
        print(f"\n🔍 Testing {filename}...")
        
        try:
            # Create parse request
            request = ParseRequest(
                content=content,
                filename=filename,
                options=ParserOptions(enable_ocr=False)  # Disable OCR for speed
            )
            
            # Parse with dispatcher
            start_time = time.time()
            result = parse_sync_dispatcher(request)
            parse_time = time.time() - start_time
            
            # Store results
            results[filename] = {
                'success': True,
                'parse_time': parse_time,
                'pages': len(result.pages),
                'tables': len(result.tables),
                'char_count': result.metadata.char_count,
                'page_count': result.metadata.page_count
            }
            
            print(f"  ✅ Success: {parse_time:.2f}s, {len(result.pages)} pages, {result.metadata.char_count} chars")
            
        except Exception as e:
            error_msg = str(e)
            errors[filename] = error_msg
            
            # Check for specific error types
            if "pandas and openpyxl are required" in error_msg:
                print(f"  ⚠️  Missing dependency: pandas/openpyxl")
            elif "beautifulsoup4 is required" in error_msg:
                print(f"  ⚠️  Missing dependency: beautifulsoup4")
            elif "python-docx is required" in error_msg:
                print(f"  ⚠️  Missing dependency: python-docx")
            elif "python-pptx is required" in error_msg:
                print(f"  ⚠️  Missing dependency: python-pptx")
            elif "pytesseract and Pillow are required" in error_msg:
                print(f"  ⚠️  Missing dependency: pytesseract/Pillow")
            elif "Unsupported format" in error_msg:
                print(f"  ❌ Unsupported format: {error_msg}")
            elif "File size" in error_msg and "exceeds limit" in error_msg:
                print(f"  ❌ File too large: {error_msg}")
            else:
                print(f"  ❌ Error: {error_msg}")
    
    # Summary
    print("\n" + "=" * 50)
    print("📊 TEST SUMMARY")
    print("=" * 50)
    
    if results:
        print(f"\n✅ Successful parses: {len(results)}")
        for filename, result in results.items():
            print(f"  {filename}: {result['parse_time']:.2f}s, {result['pages']} pages, {result['char_count']} chars")
    
    if errors:
        print(f"\n❌ Failed parses: {len(errors)}")
        for filename, error in errors.items():
            print(f"  {filename}: {error}")
    
    # Test unsupported format
    print(f"\n🔍 Testing unsupported format...")
    try:
        unsupported_content = b"This is not a recognized format"
        request = ParseRequest(
            content=unsupported_content,
            filename="test.unknown",
            options=ParserOptions()
        )
        result = parse_sync_dispatcher(request)
        print("  ⚠️  Unexpected success for unsupported format")
    except Exception as e:
        if "Unsupported format" in str(e):
            print("  ✅ Correctly rejected unsupported format")
        else:
            print(f"  ❌ Unexpected error: {e}")
    
    # Test large file
    print(f"\n🔍 Testing file size limit...")
    try:
        large_content = b"x" * (51 * 1024 * 1024)  # 51MB
        request = ParseRequest(
            content=large_content,
            filename="large.txt",
            options=ParserOptions(max_file_size=50 * 1024 * 1024)
        )
        result = parse_sync_dispatcher(request)
        print("  ⚠️  Unexpected success for oversized file")
    except Exception as e:
        if "File size" in str(e) and "exceeds limit" in str(e):
            print("  ✅ Correctly rejected oversized file")
        else:
            print(f"  ❌ Unexpected error: {e}")
    
    print(f"\n🎉 Dispatcher testing completed!")


def test_ocr_functionality():
    """Test OCR functionality if available."""
    print("\n🔍 Testing OCR functionality...")
    
    try:
        from PIL import Image, ImageDraw
        import pytesseract
        
        # Create test image with text
        img = Image.new('RGB', (400, 200), color='white')
        draw = ImageDraw.Draw(img)
        draw.text((50, 50), "OCR Test", fill='black')
        draw.text((50, 100), "This text should be extracted", fill='black')
        
        buffer = io.BytesIO()
        img.save(buffer, format='JPEG')
        content = buffer.getvalue()
        
        # Test with OCR enabled
        request = ParseRequest(
            content=content,
            filename="ocr_test.jpg",
            options=ParserOptions(enable_ocr=True)
        )
        
        result = parse_sync_dispatcher(request)
        
        if result.metadata.char_count > 0:
            print("  ✅ OCR working: extracted text successfully")
            print(f"  📝 Extracted text: {result.pages[0][:100]}...")
        else:
            print("  ⚠️  OCR enabled but no text extracted")
            
    except ImportError:
        print("  ⚠️  OCR dependencies not available (pytesseract/Pillow)")
    except Exception as e:
        print(f"  ❌ OCR test failed: {e}")


if __name__ == "__main__":
    test_dispatcher()
    test_ocr_functionality() 