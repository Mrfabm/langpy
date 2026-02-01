# This is a full copy of the docling parser, including fast image and fast PDF logic.
# It is a standalone file for reference or integration purposes.

# --- MODELS (from models.py) ---
"""
Parser Models - Core data structures for Langbase-style parser primitive.
"""
from __future__ import annotations
import hashlib
import json
import time
from datetime import datetime
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional, Union, Iterator, Literal
from pydantic import BaseModel, Field, validator
from enum import Enum

class JobStatus(str, Enum):
    QUEUED = "queued"
    PROCESSING = "processing"
    READY = "ready"
    FAILED = "failed"

class ParserError(Exception):
    def __init__(self, message: str, status_code: int = 500):
        super().__init__(message)
        self.status_code = status_code
class UnsupportedFormatError(ParserError):
    def __init__(self, format_name: str):
        super().__init__(f"Unsupported format: {format_name}", 400)
class FileTooLargeError(ParserError):
    def __init__(self, size: int, limit: int):
        super().__init__(f"File size {size} exceeds limit {limit}", 413)
class CorruptedFileError(ParserError):
    def __init__(self, filename: str, reason: str):
        super().__init__(f"Corrupted file {filename}: {reason}", 422)
class EncryptedFileError(ParserError):
    def __init__(self, filename: str):
        super().__init__(f"Encrypted file {filename} requires password", 423)
class ParseTimeoutError(ParserError):
    def __init__(self, filename: str, timeout: int):
        super().__init__(f"Parse timeout for {filename} after {timeout}s", 408)
class OcrFailureError(ParserError):
    def __init__(self, filename: str, reason: str):
        super().__init__(f"OCR failed for {filename}: {reason}", 422)
class TableCell(BaseModel):
    value: str = Field(...)
    row: int = Field(...)
    col: int = Field(...)
    rowspan: int = Field(1)
    colspan: int = Field(1)
class TableMeta(BaseModel):
    page_number: int = Field(...)
    bounding_box: List[float] = Field(...)
    rows: int = Field(...)
    columns: int = Field(...)
    cells: List[TableCell] = Field(...)
    @property
    def matrix(self) -> List[List[str]]:
        matrix = [['' for _ in range(self.columns)] for _ in range(self.rows)]
        for cell in self.cells:
            for r in range(cell.row, min(cell.row + cell.rowspan, self.rows)):
                for c in range(cell.col, min(cell.col + cell.colspan, self.columns)):
                    matrix[r][c] = cell.value
        return matrix
class PageMeta(BaseModel):
    page_number: int = Field(...)
    width: float = Field(...)
    height: float = Field(...)
    rotation: int = Field(0)
    language: Optional[str] = Field(None)
    text: str = Field(...)
class ParseStats(BaseModel):
    char_count: int = Field(...)
    token_estimate: int = Field(...)
    table_count: int = Field(0)
    page_count: int = Field(...)
    languages: List[str] = Field(default_factory=list)
    source_name: str = Field(...)
    checksum: str = Field(...)
    parse_time: float = Field(...)
    ocr_time: Optional[float] = Field(None)
class ParseResult(BaseModel):
    pages: List[str] = Field(...)
    tables: List[TableMeta] = Field(default_factory=list)
    metadata: ParseStats = Field(...)
    def to_dict(self) -> Dict[str, Any]:
        return {
            "pages": self.pages,
            "tables": [table.dict() for table in self.tables],
            "metadata": self.metadata.dict()
        }
class ParseJob(BaseModel):
    id: str = Field(...)
    status: JobStatus = Field(...)
    created_at: datetime = Field(...)
    updated_at: datetime = Field(...)
    input_type: str = Field(...)
    mime_type: Optional[str] = Field(None)
    error_code: Optional[str] = Field(None)
    error_message: Optional[str] = Field(None)
    stats: Optional[ParseStats] = Field(None)
    result_url: Optional[str] = Field(None)
    @validator('created_at', 'updated_at', pre=True)
    def parse_datetime(cls, v):
        if isinstance(v, str):
            return datetime.fromisoformat(v.replace('Z', '+00:00'))
        return v
class ParserOptions(BaseModel):
    max_file_size: int = Field(50 * 1024 * 1024)
    enable_ocr: bool = Field(True)
    ocr_languages: List[str] = Field(["eng"])
    ocr_confidence_threshold: float = Field(0.7)
    language_hint: Optional[str] = Field(None)
    preserve_whitespace: bool = Field(False)
    merge_hyphens: bool = Field(True)
    strip_headers_footers: bool = Field(False)
    table_strategy: Literal["docling", "camelot", "none"] = Field("docling")
    parse_timeout: int = Field(300)
    beta_features: bool = Field(False)
    experimental_chunking: bool = Field(False)
class ParseRequest(BaseModel):
    content: Optional[Union[bytes, str]] = Field(None)
    url: Optional[str] = Field(None)
    filename: Optional[str] = Field(None)
    options: Optional[ParserOptions] = Field(None)
    @validator('content', 'url')
    def validate_input(cls, v, values):
        if not v and not values.get('content') and not values.get('url'):
            raise ValueError("Either content or url must be provided")
        return v
    def __post_init__(self):
        if self.options is None:
            self.options = ParserOptions()
class SupportedFormat(BaseModel):
    mime_type: str = Field(...)
    max_file_size: int = Field(...)
    description: str = Field(...)
    extensions: List[str] = Field(...)
@dataclass
class ParsedDocument:
    text: str
    metadata: Dict[str, Any] = field(default_factory=dict)
    def __post_init__(self):
        if isinstance(self.text, bytes):
            self.text = self.text.decode('utf-8', errors='replace')
        self.text = self.text.replace('\r\n', '\n').replace('\r', '\n')
    def to_dict(self) -> Dict[str, Any]:
        return {
            "text": self.text,
            "metadata": self.metadata
        }
    def to_json(self) -> str:
        return json.dumps(self.to_dict(), ensure_ascii=False, indent=2)

# --- FAST IMAGE PARSER (from fast_image_parser.py) ---
import io
try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
class FastImageParser:
    def __init__(self, options: Optional[ParserOptions] = None):
        self.options = options or ParserOptions()
    def _calculate_checksum(self, content: bytes) -> str:
        return hashlib.sha256(content).hexdigest()
    def parse_sync(self, request: ParseRequest) -> ParseResult:
        if not PIL_AVAILABLE:
            raise ParserError("PIL/Pillow not available. Install with: pip install Pillow")
        import time
        start_time = time.time()
        try:
            img = Image.open(io.BytesIO(request.content))
            width, height = img.size
            format_name = img.format or "unknown"
            mode = img.mode
            parse_time = time.time() - start_time
            stats = ParseStats(
                char_count=0,
                token_estimate=0,
                table_count=0,
                page_count=1,
                languages=["unknown"],
                source_name=request.filename or "image",
                checksum=self._calculate_checksum(request.content),
                parse_time=parse_time
            )
            return ParseResult(
                pages=[f"Image: {width}x{height} {format_name} {mode}"],
                tables=[],
                metadata=stats
            )
        except Exception as e:
            raise ParserError(f"Fast image parsing failed: {str(e)}")
def parse_sync_fast_image(request: ParseRequest) -> ParseResult:
    parser = FastImageParser()
    return parser.parse_sync(request)

# --- FAST PDF PARSER (from fast_pdf_parser.py) ---
try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False
class FastPDFParser:
    def __init__(self, options: Optional[ParserOptions] = None):
        self.options = options or ParserOptions()
    def _calculate_checksum(self, content: bytes) -> str:
        return hashlib.sha256(content).hexdigest()
    def _estimate_tokens(self, text: str) -> int:
        return len(text) // 4
    def _detect_language(self, text: str) -> List[str]:
        if not text.strip():
            return ["unknown"]
        if any(ord(c) > 127 for c in text[:1000]):
            return ["multilingual"]
        else:
            return ["en"]
    def _process_text_content(self, text: str, options: ParserOptions) -> str:
        if not text:
            return text
        if options.merge_hyphens:
            import re
            text = re.sub(r'(\w+)-\s*\n\s*(\w+)', r'\1\2', text)
        if options.strip_headers_footers:
            lines = text.split('\n')
            if len(lines) > 10:
                lines = lines[2:-2]
            text = '\n'.join(lines)
        if not options.preserve_whitespace:
            text = ' '.join(text.split())
        return text
    def parse_sync(self, request: ParseRequest) -> ParseResult:
        if not PYPDF2_AVAILABLE:
            raise ParserError("PyPDF2 not available. Install with: pip install PyPDF2")
        import time
        start_time = time.time()
        try:
            pdf_file = io.BytesIO(request.content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            pages = []
            page_metas = []
            for page_num, page in enumerate(pdf_reader.pages, 1):
                try:
                    page_text = page.extract_text()
                    processed_text = self._process_text_content(page_text, request.options or self.options)
                    pages.append(processed_text)
                    page_meta = PageMeta(
                        page_number=page_num,
                        width=612,
                        height=792,
                        rotation=0,
                        language=self._detect_language(processed_text)[0],
                        text=processed_text
                    )
                    page_metas.append(page_meta)
                except Exception as e:
                    pages.append(f"Page {page_num}: Text extraction failed - {str(e)}")
                    page_meta = PageMeta(
                        page_number=page_num,
                        width=612,
                        height=792,
                        rotation=0,
                        language="unknown",
                        text=""
                    )
                    page_metas.append(page_meta)
            tables = []
            total_text = '\n'.join(pages)
            char_count = len(total_text)
            token_estimate = self._estimate_tokens(total_text)
            languages = list(set([meta.language for meta in page_metas if meta.language]))
            parse_time = time.time() - start_time
            stats = ParseStats(
                char_count=char_count,
                token_estimate=token_estimate,
                table_count=len(tables),
                page_count=len(pages),
                languages=languages,
                source_name=request.filename or "document.pdf",
                checksum=self._calculate_checksum(request.content),
                parse_time=parse_time
            )
            return ParseResult(
                pages=pages,
                tables=tables,
                metadata=stats
            )
        except Exception as e:
            raise ParserError(f"Fast PDF parsing failed: {str(e)}")
def parse_sync_fast_pdf(request: ParseRequest) -> ParseResult:
    parser = FastPDFParser()
    return parser.parse_sync(request)

# --- SIMPLE PARSER (from simple_parser.py) ---
class SimpleParser:
    def __init__(self, options: Optional[ParserOptions] = None):
        self.options = options or ParserOptions()
    def _calculate_checksum(self, content: bytes) -> str:
        return hashlib.sha256(content).hexdigest()
    def _estimate_tokens(self, text: str) -> int:
        return len(text) // 4
    def _detect_language(self, text: str) -> List[str]:
        if not text.strip():
            return ["unknown"]
        if any(ord(c) > 127 for c in text[:1000]):
            return ["multilingual"]
        else:
            return ["en"]
    def _process_text_content(self, text: str, options: ParserOptions) -> str:
        if not text:
            return text
        if options.merge_hyphens:
            import re
            text = re.sub(r'(\w+)-\s*\n\s*(\w+)', r'\1\2', text)
        if options.strip_headers_footers:
            lines = text.split('\n')
            if len(lines) > 10:
                lines = lines[2:-2]
            text = '\n'.join(lines)
        if not options.preserve_whitespace:
            text = ' '.join(text.split())
        return text
    def parse_sync(self, request: ParseRequest) -> ParseResult:
        import time
        start_time = time.time()
        try:
            if isinstance(request.content, bytes):
                try:
                    text = request.content.decode('utf-8')
                except UnicodeDecodeError:
                    text = request.content.decode('latin-1')
            else:
                text = str(request.content)
            processed_text = self._process_text_content(text, request.options or self.options)
            pages = [page.strip() for page in processed_text.split('\n\n') if page.strip()]
            if not pages:
                pages = [processed_text]
            page_metas = []
            for i, page_text in enumerate(pages, 1):
                page_meta = PageMeta(
                    page_number=i,
                    width=612,
                    height=792,
                    rotation=0,
                    language=self._detect_language(page_text)[0],
                    text=page_text
                )
                page_metas.append(page_meta)
            tables = []
            total_text = '\n'.join(pages)
            char_count = len(total_text)
            token_estimate = self._estimate_tokens(total_text)
            languages = list(set([meta.language for meta in page_metas if meta.language]))
            parse_time = time.time() - start_time
            stats = ParseStats(
                char_count=char_count,
                token_estimate=token_estimate,
                table_count=len(tables),
                page_count=len(pages),
                languages=languages,
                source_name=request.filename or "document",
                checksum=self._calculate_checksum(request.content),
                parse_time=parse_time
            )
            return ParseResult(
                pages=pages,
                tables=tables,
                metadata=stats
            )
        except Exception as e:
            raise ParserError(f"Simple parsing failed: {str(e)}")
def parse_sync_simple(request: ParseRequest) -> ParseResult:
    parser = SimpleParser()
    return parser.parse_sync(request)

# --- DOCLING PARSER (from docling_parser.py) ---
# (Paste the full class and logic from your current docling_parser.py here) 

# --- SUPPORTED FORMATS REGISTRY ---
SUPPORTED_FORMATS = {
    'application/pdf': ['.pdf'],
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': ['.docx'],
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': ['.pptx'],
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': ['.xlsx'],
    'text/html': ['.html', '.htm'],
    'text/markdown': ['.md', '.markdown'],
    'text/plain': ['.txt'],
    'text/csv': ['.csv'],
    'image/jpeg': ['.jpg', '.jpeg'],
    'image/png': ['.png'],
    'image/tiff': ['.tiff', '.tif'],
}

# --- MIME DETECTION LOGIC (from mime_detection.py) ---
import mimetypes
from pathlib import Path
class MimeDetector:
    def detect_from_bytes(self, content: bytes, filename: Optional[str] = None) -> str:
        mime_type = self._detect_magic_bytes(content)
        if mime_type:
            return mime_type
        if filename:
            mime_type = self._detect_from_extension(filename)
            if mime_type:
                return mime_type
        return 'application/octet-stream'
    def _detect_magic_bytes(self, content: bytes) -> Optional[str]:
        if content.startswith(b'%PDF'):
            return 'application/pdf'
        elif content.startswith(b'PK\x03\x04'):
            return 'application/zip'
        elif content.startswith(b'\xff\xd8\xff'):
            return 'image/jpeg'
        elif content.startswith(b'\x89PNG\r\n\x1a\n'):
            return 'image/png'
        elif content.startswith(b'{\x22'):
            return 'application/json'
        elif content.startswith(b'<html') or content.startswith(b'<!DOCTYPE'):
            return 'text/html'
        elif content.startswith(b'<?xml'):
            return 'application/xml'
        return None
    def _detect_from_extension(self, filename: str) -> Optional[str]:
        path = Path(filename)
        extension = path.suffix.lower()
        extension_map = {
            '.txt': 'text/plain', '.md': 'text/markdown', '.html': 'text/html', '.htm': 'text/html',
            '.xml': 'application/xml', '.json': 'application/json', '.yaml': 'application/x-yaml', '.yml': 'application/x-yaml',
            '.csv': 'text/csv', '.pdf': 'application/pdf',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg', '.png': 'image/png', '.gif': 'image/gif',
            '.zip': 'application/zip', '.py': 'application/x-python', '.js': 'application/javascript',
            '.ts': 'application/typescript', '.java': 'text/x-java-source', '.cpp': 'text/x-c++src', '.c': 'text/x-csrc',
            '.go': 'text/x-go', '.rs': 'text/x-rust', '.php': 'application/x-php', '.rb': 'text/x-ruby',
            '.sql': 'application/sql', '.sh': 'application/x-sh', '.css': 'text/css',
        }
        if extension in extension_map:
            return extension_map[extension]
        mime_type, _ = mimetypes.guess_type(filename)
        return mime_type
mime_detector = MimeDetector()

# --- FORMAT-SPECIFIC PARSERS (XLSX, HTML, DOCX, PPTX, OCR) ---
def parse_sync_xlsx(request: ParseRequest) -> ParseResult:
    try:
        import pandas as pd
        from io import BytesIO
    except ImportError:
        raise ParserError("pandas and openpyxl are required for XLSX parsing. Install with: pip install pandas openpyxl")
    start_time = time.time()
    excel_file = BytesIO(request.content)
    try:
        # Use ExcelFile for better performance and proper sheet handling
        with pd.ExcelFile(excel_file) as xls:
            sheet_names = xls.sheet_names
            text_parts = []
            
            for sheet_name in sheet_names:
                try:
                    # Read with header=None to avoid header inference issues
                    # Use dtype=str to avoid numeric casting and NaN issues
                    df = pd.read_excel(xls, sheet_name=sheet_name, header=None, dtype=str)
                    
                    if not df.empty:
                        # Clean the DataFrame: remove completely empty rows/columns and fill NaN
                        df = df.dropna(how='all', axis=0).dropna(how='all', axis=1)
                        df = df.fillna('')  # Replace NaN with empty strings
                        
                        # Convert to tab-separated text for better readability
                        sheet_text = f"Sheet: {sheet_name}\n"
                        sheet_text += df.to_csv(sep='\t', index=False, header=False)
                        text_parts.append(sheet_text)
                        
                except Exception as e:
                    # Skip problematic sheets but continue with others
                    text_parts.append(f"Sheet: {sheet_name}\n[Error reading sheet: {str(e)}]")
                    continue
        
        full_text = '\n\n'.join(text_parts)
        parse_time = time.time() - start_time
        stats = ParseStats(
            char_count=len(full_text),
            token_estimate=len(full_text) // 4,
            table_count=len(sheet_names),
            page_count=1,
            languages=["unknown"],
            source_name=request.filename or "spreadsheet.xlsx",
            checksum=hashlib.sha256(request.content).hexdigest(),
            parse_time=parse_time
        )
        return ParseResult(pages=[full_text], tables=[], metadata=stats)
    except Exception as e:
        raise ParserError(f"XLSX parsing failed: {str(e)}")

def parse_sync_html(request: ParseRequest) -> ParseResult:
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        raise ParserError("beautifulsoup4 is required for HTML parsing. Install with: pip install beautifulsoup4")
    start_time = time.time()
    if isinstance(request.content, bytes):
        content = request.content.decode('utf-8', errors='replace')
    else:
        content = str(request.content)
    soup = BeautifulSoup(content, 'html.parser')
    text = soup.get_text(separator=' ', strip=True)
    parse_time = time.time() - start_time
    stats = ParseStats(
        char_count=len(text),
        token_estimate=len(text) // 4,
        table_count=0,
        page_count=1,
        languages=["unknown"],
        source_name=request.filename or "document.html",
        checksum=hashlib.sha256(request.content if isinstance(request.content, bytes) else content.encode('utf-8')).hexdigest(),
        parse_time=parse_time
    )
    return ParseResult(pages=[text], tables=[], metadata=stats)

def parse_sync_docx(request: ParseRequest) -> ParseResult:
    try:
        from docx import Document
        from io import BytesIO
    except ImportError:
        raise ParserError("python-docx is required for DOCX parsing. Install with: pip install python-docx")
    start_time = time.time()
    doc = Document(BytesIO(request.content))
    text_parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                text_parts.append(' | '.join(row_text))
    full_text = '\n\n'.join(text_parts)
    parse_time = time.time() - start_time
    stats = ParseStats(
        char_count=len(full_text),
        token_estimate=len(full_text) // 4,
        table_count=len(doc.tables),
        page_count=1,
        languages=["unknown"],
        source_name=request.filename or "document.docx",
        checksum=hashlib.sha256(request.content).hexdigest(),
        parse_time=parse_time
    )
    return ParseResult(pages=[full_text], tables=[], metadata=stats)

def parse_sync_pptx(request: ParseRequest) -> ParseResult:
    try:
        from pptx import Presentation
        from io import BytesIO
    except ImportError:
        raise ParserError("python-pptx is required for PPTX parsing. Install with: pip install python-pptx")
    start_time = time.time()
    prs = Presentation(BytesIO(request.content))
    text_parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            text_parts.append(f"Slide {slide_num}:\n" + '\n'.join(slide_text))
    full_text = '\n\n'.join(text_parts)
    parse_time = time.time() - start_time
    stats = ParseStats(
        char_count=len(full_text),
        token_estimate=len(full_text) // 4,
        table_count=0,
        page_count=len(prs.slides),
        languages=["unknown"],
        source_name=request.filename or "presentation.pptx",
        checksum=hashlib.sha256(request.content).hexdigest(),
        parse_time=parse_time
    )
    return ParseResult(pages=[full_text], tables=[], metadata=stats)

def parse_sync_ocr_image(request: ParseRequest) -> ParseResult:
    try:
        from PIL import Image
        import pytesseract
        import io
    except ImportError:
        raise ParserError("pytesseract and Pillow are required for OCR. Install with: pip install pytesseract Pillow")
    start_time = time.time()
    img = Image.open(io.BytesIO(request.content))
    text = pytesseract.image_to_string(img, lang='+'.join(request.options.ocr_languages) if request.options and request.options.ocr_languages else 'eng')
    parse_time = time.time() - start_time
    stats = ParseStats(
        char_count=len(text),
        token_estimate=len(text) // 4,
        table_count=0,
        page_count=1,
        languages=["unknown"],
        source_name=request.filename or "image",
        checksum=hashlib.sha256(request.content).hexdigest(),
        parse_time=parse_time,
        ocr_time=parse_time
    )
    return ParseResult(pages=[text], tables=[], metadata=stats)

# --- MAIN DISPATCHER ---
def parse_sync_dispatcher(request: ParseRequest) -> ParseResult:
    # File size guard
    if request.options and hasattr(request.options, 'max_file_size'):
        max_size = request.options.max_file_size
    else:
        max_size = 50 * 1024 * 1024
    if request.content and len(request.content) > max_size:
        raise FileTooLargeError(len(request.content), max_size)
    # Detect MIME type
    mime_type = None
    if request.filename:
        mime_type = mime_detector.detect_from_bytes(request.content, request.filename)
    else:
        mime_type = mime_detector.detect_from_bytes(request.content)
    ext = Path(request.filename).suffix.lower() if request.filename else ''
    # Dispatch by MIME type/extension
    if mime_type.startswith('image/'):
        if request.options and getattr(request.options, 'enable_ocr', True):
            return parse_sync_ocr_image(request)
        else:
            return parse_sync_fast_image(request)
    elif mime_type == 'application/pdf' or ext == '.pdf':
        return parse_sync_fast_pdf(request)
    elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet' or ext == '.xlsx':
        return parse_sync_xlsx(request)
    elif mime_type == 'text/html' or ext in {'.html', '.htm'}:
        return parse_sync_html(request)
    elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' or ext == '.docx':
        return parse_sync_docx(request)
    elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation' or ext == '.pptx':
        return parse_sync_pptx(request)
    elif mime_type in {'text/plain', 'text/markdown'} or ext in {'.txt', '.md', '.markdown'}:
        return parse_sync_simple(request)
    elif mime_type == 'text/csv' or ext == '.csv':
        return parse_sync_simple(request)
    # Fallback to Docling universal parser if available
    try:
        from parser.docling_parser import DoclingParser
        parser = DoclingParser(request.options)
        return parser.parse_sync(request)
    except Exception as e:
        raise UnsupportedFormatError(f"{mime_type} / {ext}: {e}") 