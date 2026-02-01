"""
Office Parser - Basic office document text extraction.

Handles DOCX, PPTX, and XLSX documents with simple text extraction.
"""

import time
from typing import Dict, Any, Optional
from ..models import ParsedDocument, ParseRequest, ParseResult, ParserError


class OfficeParser:
    """Parser for office documents (DOCX, PPTX, XLSX)."""
    
    def __init__(self):
        """Initialize office parser."""
        pass
    
    async def parse(self, request: ParseRequest) -> ParseResult:
        """
        Parse office document.
        
        Args:
            request: Parse request with office document content
            
        Returns:
            Parse result with extracted text
        """
        start_time = time.time()
        
        try:
            # Get content as bytes
            if isinstance(request.content, str):
                content = request.content.encode('utf-8')
            else:
                content = request.content
            
            # Extract text based on content type
            text, metadata = await self._extract_text(content, request)
            
            # Create document
            document = ParsedDocument(text=text, metadata=metadata)
            
            parse_time = time.time() - start_time
            return ParseResult(document=document, parse_time=parse_time)
            
        except Exception as e:
            raise ParserError(f"Office document parsing failed: {str(e)}")
    
    async def _extract_text(self, content: bytes, request: ParseRequest) -> tuple[str, Dict[str, Any]]:
        """Extract text from office document content."""
        content_type = request.content_type or 'application/zip'
        
        if 'wordprocessingml.document' in content_type:
            return await self._extract_docx(content, request)
        elif 'presentationml.presentation' in content_type:
            return await self._extract_pptx(content, request)
        elif 'spreadsheetml.sheet' in content_type:
            return await self._extract_xlsx(content, request)
        else:
            raise ParserError(f"Unsupported office format: {content_type}")
    
    async def _extract_docx(self, content: bytes, request: ParseRequest) -> tuple[str, Dict[str, Any]]:
        """Extract text from DOCX document."""
        try:
            from docx import Document
            from io import BytesIO
            
            # Load document
            doc = Document(BytesIO(content))
            
            # Extract text from paragraphs
            text_parts = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_parts.append(' | '.join(row_text))
            
            full_text = '\n\n'.join(text_parts)
            
            metadata = {
                'content_type': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                'filename': request.filename,
                'text_length': len(full_text),
                'paragraph_count': len([p for p in doc.paragraphs if p.text.strip()]),
                'table_count': len(doc.tables)
            }
            
            return full_text, metadata
            
        except Exception as e:
            raise ParserError(f"DOCX extraction failed: {str(e)}")
    
    async def _extract_pptx(self, content: bytes, request: ParseRequest) -> tuple[str, Dict[str, Any]]:
        """Extract text from PPTX document."""
        try:
            from pptx import Presentation
            from io import BytesIO
            
            # Load presentation
            prs = Presentation(BytesIO(content))
            
            # Extract text from slides
            text_parts = []
            slide_count = len(prs.slides)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    text_parts.append(f"Slide {slide_num}:\n" + '\n'.join(slide_text))
            
            full_text = '\n\n'.join(text_parts)
            
            metadata = {
                'content_type': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                'filename': request.filename,
                'text_length': len(full_text),
                'slide_count': slide_count
            }
            
            return full_text, metadata
            
        except Exception as e:
            raise ParserError(f"PPTX extraction failed: {str(e)}")
    
    async def _extract_xlsx(self, content: bytes, request: ParseRequest) -> tuple[str, Dict[str, Any]]:
        """Extract text from XLSX document."""
        try:
            import pandas as pd
            from io import BytesIO
            
            # Load workbook
            excel_file = BytesIO(content)
            
            # Read all sheets
            text_parts = []
            sheet_names = pd.ExcelFile(excel_file).sheet_names
            
            for sheet_name in sheet_names:
                try:
                    df = pd.read_excel(excel_file, sheet_name=sheet_name)
                    
                    if not df.empty:
                        # Convert DataFrame to text
                        sheet_text = f"Sheet: {sheet_name}\n"
                        sheet_text += df.to_string(index=False)
                        text_parts.append(sheet_text)
                        
                except Exception:
                    # Skip problematic sheets
                    continue
            
            full_text = '\n\n'.join(text_parts)
            
            metadata = {
                'content_type': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                'filename': request.filename,
                'text_length': len(full_text),
                'sheet_count': len(sheet_names)
            }
            
            return full_text, metadata
            
        except Exception as e:
            raise ParserError(f"XLSX extraction failed: {str(e)}") 